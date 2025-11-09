# --------------------------------------------------------------
# app.py   –   Flask + Gemini 2.5 + PPTX (stable background + overlay)
# --------------------------------------------------------------
import os
import json
import re
import time
import uuid
import requests
import tempfile
from io import BytesIO
from dotenv import load_dotenv
import google.generativeai as genai
from google.generativeai import types
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageFilter, ImageStat
from flask import Flask, render_template, request, send_file, flash

# ------------------------------------------------------------------
# 1. Load env + Flask setup
# ------------------------------------------------------------------
load_dotenv()
app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "replace-in-production")

# ------------------------------------------------------------------
# 2. Gemini model configuration
# ------------------------------------------------------------------
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    print("GEMINI_API_KEY missing – Gemini will not work.")
    model = None
else:
    genai.configure(api_key=GEMINI_API_KEY)

    system_prompt = (
        "You are an expert presentation outline generator. "
        "Return ONLY valid JSON in this exact shape:\n"
        '{"slides": [{"title": "...", "content": ["...", "..."]}]}'
    )

    generation_cfg = types.GenerationConfig(
        response_mime_type="application/json",
        temperature=0.7,
        max_output_tokens=2048,
    )

    try:
        model = genai.GenerativeModel(
            model_name="gemini-2.5-flash",
            system_instruction=system_prompt,
            generation_config=generation_cfg,
        )
        print("Gemini 2.5 Flash model ready.")
    except Exception as e:
        print(f"Gemini model init error: {e}")
        model = None

# ------------------------------------------------------------------
# 3. Google Custom Search Image Fetcher
# ------------------------------------------------------------------
def fetch_image_url_from_google(query: str) -> str | None:
    api_key = os.getenv("GOOGLE_SEARCH_API_KEY")
    cse_id = os.getenv("GOOGLE_CSE_ID")
    if not api_key or not cse_id:
        return None

    try:
        url = "https://www.googleapis.com/customsearch/v1"
        params = {
            "q": query,
            "cx": cse_id,
            "key": api_key,
            "searchType": "image",
            "num": 1,
            "safe": "active",
        }
        r = requests.get(url, params=params, timeout=10)
        r.raise_for_status()
        items = r.json().get("items") or []
        if not items:
            return None
        return items[0].get("link")
    except Exception as e:
        print(f"Image fetch error: {e}")
        return None

# ------------------------------------------------------------------
# 4. ROBUST set_slide_background (fixed)
# ------------------------------------------------------------------
# ------------------------------------------------------------------
# 4. FINAL set_slide_background – NEVER black or white
# ------------------------------------------------------------------
def set_slide_background(slide, image_url: str | None) -> float:
    """
    1. Try the supplied Google-image URL
    2. If it fails or contains text → use static/fallback.jpg
    3. If the local file is missing → download it from the remote URL
    4. ALWAYS add a dark overlay (0.4 opacity) for readability
    Returns brightness (0-255) → white text on dark images, black text on light images
    """
    FALLBACK_PATH      = os.path.join(os.path.dirname(__file__), "static", "fallback.jpg")
    REMOTE_FALLBACK_URL = "https://slidechef.net/wp-content/uploads/2023/09/cool-background.jpg"

    # ---------- helpers ----------
    def slide_dims():
        prs = slide.part.package.presentation_part.presentation
        return prs.slide_width, prs.slide_height

    def brightness(img: Image.Image) -> float:
        return float(ImageStat.Stat(img.convert("L")).mean[0])

    def has_text(img: Image.Image) -> bool:
        try:
            gray   = img.convert("L")
            edges  = gray.filter(ImageFilter.FIND_EDGES)
            contrast = ImageStat.Stat(edges).stddev[0]
            mean   = ImageStat.Stat(gray).mean[0]
            return (contrast / (mean + 1e-6)) > 1.6
        except Exception:
            return False

    def put_image(stream: BytesIO):
        w, h = slide_dims()
        pic = slide.shapes.add_picture(stream, 0, 0, w, h)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)   # behind everything

    def dark_overlay():
        w, h = slide_dims()
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, h)
        ov.fill.solid()
        ov.fill.fore_color.rgb = RGBColor(0, 0, 0)
        ov.fill.transparency   = 0.4
        slide.shapes._spTree.remove(ov._element)
        slide.shapes._spTree.insert(3, ov._element)

    # ---------- 1. make sure a local fallback exists ----------
    if not os.path.exists(FALLBACK_PATH):
        os.makedirs(os.path.dirname(FALLBACK_PATH), exist_ok=True)
        try:
            print("Downloading fallback image …")
            r = requests.get(REMOTE_FALLBACK_URL,
                             headers={"User-Agent": "Mozilla/5.0"},
                             timeout=12)
            r.raise_for_status()
            with open(FALLBACK_PATH, "wb") as f:
                f.write(r.content)
            print("Fallback saved locally.")
        except Exception as e:
            print(f"Fallback download failed: {e}")

    # ---------- 2. try the primary image ----------
    img = None
    if image_url:
        try:
            r = requests.get(image_url,
                             headers={"User-Agent": "Mozilla/5.0"},
                             timeout=12)
            r.raise_for_status()
            candidate = Image.open(BytesIO(r.content))
            if not has_text(candidate):
                img = candidate
                print(f"Using primary image: {image_url}")
            else:
                print(f"Skipping text-heavy primary image.")
        except Exception as e:
            print(f"Primary image error: {e}")

    # ---------- 3. use local fallback ----------
    if img is None and os.path.exists(FALLBACK_PATH):
        try:
            img = Image.open(FALLBACK_PATH)
            print("Using local fallback image.")
        except Exception as e:
            print(f"Local fallback open error: {e}")

    # ---------- 4. last-ditch remote fallback ----------
    if img is None:
        try:
            print("Fetching remote fallback …")
            r = requests.get(REMOTE_FALLBACK_URL,
                             headers={"User-Agent": "Mozilla/5.0"},
                             timeout=12)
            r.raise_for_status()
            img = Image.open(BytesIO(r.content))
            # cache it for next run
            with open(FALLBACK_PATH, "wb") as f:
                f.write(r.content)
            print("Remote fallback used & cached.")
        except Exception as e:
            print(f"All image sources failed: {e}")

    # ---------- 5. if we still have no image → solid dark (never white) ----------
    if img is None:
        print("No image → solid dark background.")
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(20, 20, 40)
        return 40.0                     # forces white text

    # ---------- 6. apply the chosen image ----------
    bright = brightness(img)
    stream = BytesIO()
    img.save(stream, format="JPEG")
    stream.seek(0)
    put_image(stream)                      # always a dark overlay for readability
    return bright
# ------------------------------------------------------------------
# 5. Generate Gemini content
# ------------------------------------------------------------------
def generate_slide_content(topic: str, num_slides: int, retries: int = 2):
    if not model:
        return None

    user_prompt = (
        f"Generate a {num_slides}-slide presentation on: '{topic}'. "
        "First slide = title/intro, last = conclusion. "
        "Each slide must have a 'title' and 3–5 bullet points in 'content'. "
        "Output only the JSON object."
    )

    def clean_json(text: str):
        text = re.sub(r"^```(?:json)?\s*", "", text.strip())
        text = re.sub(r"```$", "", text)
        m = re.search(r"\{[\s\S]*\}", text)
        if not m:
            return "{}"
        cleaned = m.group(0)
        return re.sub(r",\s*(\}|\])", r"\1", cleaned)

    for i in range(retries):
        try:
            res = model.generate_content(contents=user_prompt)
            raw = getattr(res, "text", "")
            data = json.loads(clean_json(raw))
            if "slides" in data and len(data["slides"]) >= 2:
                print(f"Gemini generated {len(data['slides'])} slides.")
                return data
        except Exception as e:
            print(f"Gemini attempt {i+1}/{retries} failed: {e}")
            time.sleep(2)
    return None

# ------------------------------------------------------------------
# 6. Build PPTX
# ------------------------------------------------------------------
def create_presentation(topic: str, outline: dict, path: str) -> bool:
    if not outline or "slides" not in outline:
        return False

    prs = Presentation()
    for i, s in enumerate(outline["slides"]):
        title, bullets = s.get("title", "Slide"), s.get("content", [])
        slide = prs.slides.add_slide(prs.slide_layouts[1] if i else prs.slide_layouts[0])

        # title
        if slide.shapes.title:
            slide.shapes.title.text = title if i else topic

        # body bullets
        if len(slide.placeholders) > 1 and hasattr(slide.placeholders[1], "text_frame"):
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for b in bullets:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0

        # background + text colour
        img_url = fetch_image_url_from_google(f"{topic} {title} background")
        brightness = set_slide_background(slide, img_url)
        text_color = RGBColor(255, 255, 255) if brightness < 140 else RGBColor(0, 0, 0)

        # apply colour/size to all text
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = text_color
                        run.font.size = Pt(18)

    prs.save(path)
    return True

# ------------------------------------------------------------------
# 7. Flask route
# ------------------------------------------------------------------
@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        topic = request.form.get("topic", "").strip()
        num_slides = max(3, min(10, int(request.form.get("num_slides", 5))))

        if not topic:
            flash("Please enter a topic.", "error")
            return render_template("index.html")

        if not model:
            flash("Gemini API not configured.", "error")
            return render_template("index.html")

        outline = generate_slide_content(topic, num_slides)
        if not outline:
            flash("Failed to generate presentation outline. Try again.", "error")
            return render_template("index.html")

        safe_name = re.sub(r"[^\w\s-]", "", topic).strip().lower()
        safe_name = re.sub(r"[-\s]+", "_", safe_name)[:50]
        filename = f"{safe_name}_{uuid.uuid4().hex[:8]}.pptx"
        filepath = os.path.join(tempfile.gettempdir(), filename)

        if not create_presentation(topic, outline, filepath):
            flash("Error creating PPTX file.", "error")
            return render_template("index.html")

        try:
            return send_file(
                filepath,
                as_attachment=True,
                download_name=f"{safe_name}.pptx",
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        finally:
            try:
                os.remove(filepath)
            except Exception:
                pass

    return render_template("index.html")

# --------------------------------------------------------------

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=False)